#!/usr/bin/env python3
"""Build the 4-slide hackathon submission deck for Aberdeen hack team 05.

Layout follows the house pattern taken from Aberdeen's real client credential
slides (see docs/proposal-layout-conventions.md):

  * left-aligned descriptive slide title with a short rule to its left
  * full-width labelled section banners, one per block of content
  * a light vertical rail under each banner, with an open "donut" marker on the
    rail for every top-level bullet
  * sub-bullets step in with a filled square; a third level uses an en dash
  * a narrower second column when a list would otherwise dominate the slide
  * slide number bottom-right

Palette stays on the deck's own colours (navy #0F2540, accent #C0642A,
Calibri) rather than the Aberdeen brand template.

Geometry is measured, not hand-placed: every bullet is its own textbox whose
height comes from an estimated wrapped-line count, and blocks stack from the
returned cursor. The verification pass at the bottom re-opens the saved file
and checks slide count, fonts, sizes, off-slide shapes, unexpected overlaps
and text overflow.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ------------------------------------------------------------------ palette --
NAVY = RGBColor(0x0F, 0x25, 0x40)        # banners, titles
INK = RGBColor(0x1A, 0x27, 0x33)         # body text (near-black, not black)
INK_SOFT = RGBColor(0x33, 0x41, 0x4F)    # sub-bullet text
SLATE = RGBColor(0x4A, 0x5A, 0x6E)       # captions, page number
ACCENT = RGBColor(0xC0, 0x64, 0x2A)      # donut rings, title rule
RAIL = RGBColor(0xD3, 0xDA, 0xE2)        # dot-rail line (light grey)
LIGHT = RGBColor(0xF4, 0xF6, 0xF9)       # screenshot placeholder ground
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
QUOTE_FONT = "Georgia"

# ----------------------------------------------------------------- geometry --
SW, SH = 13.333, 7.5
M_L = 0.62                     # content left edge
CONTENT_W = 12.09              # full content width

RULE_X, RULE_W, RULE_H = 0.62, 0.42, 0.030   # short rule left of the title
TITLE_X = 1.20                               # title starts right of the rule

BANNER_H = 0.30
LABEL_DX = 0.26                # label inset inside the banner
LABEL_SZ = 12.5

RAIL_DX = 0.28                 # rail centre, from block left
RAIL_W = 0.014
DONUT_D = 0.135
TEXT_DX = 0.52                 # level-0 text left, from block left
SUB_MARK_DX = 0.78             # level-1 square marker
SUB_TEXT_DX = 0.98             # level-1 text
L3_TEXT_DX = 1.20              # level-2 text (en-dash prefixed)
SQ = 0.075                     # square marker size

BODY_SZ = 13.0
SUB_SZ = 12.5
L3_SZ = 12.0
FURNITURE_SZ = 12.0           # kickers, page number, team line — never below 12
LINE = 1.06

GAP_ITEM = 0.070               # after a level-0 bullet
GAP_SUB = 0.055                # after a sub-bullet
GAP_BANNER = 0.10              # banner bottom to first bullet
GAP_BLOCK = 0.20               # between blocks

# --------------------------------------------------------------- text metrics --
# Calibri and Georgia are not installed on Linux build hosts, so widths come
# from the metric-similar Liberation faces scaled by a calibration ratio, and
# lines are broken with the same greedy word wrap PowerPoint uses. Falls back
# to a conservative flat average if fontTools or the fonts are unavailable.
FALLBACK_CW = {"Calibri": 0.505, "Georgia": 0.545}
_LIB = "/usr/share/fonts/truetype/liberation/"
FACE = {  # our font -> (donor file, regular/bold, calibration ratio)
    ("Calibri", False): (_LIB + "LiberationSans-Regular.ttf", 0.945),
    ("Calibri", True): (_LIB + "LiberationSans-Bold.ttf", 0.945),
    ("Georgia", False): (_LIB + "LiberationSerif-Regular.ttf", 1.120),
    ("Georgia", True): (_LIB + "LiberationSerif-Bold.ttf", 1.120),
}
_metrics_cache = {}


def _metrics(font, bold):
    key = (font, bool(bold))
    if key in _metrics_cache:
        return _metrics_cache[key]
    entry = None
    spec = FACE.get(key) or FACE.get((font, False))
    if spec:
        try:
            from fontTools.ttLib import TTFont
            tt = TTFont(spec[0], lazy=True)
            upem = tt["head"].unitsPerEm
            cmap = tt.getBestCmap()
            hmtx = tt["hmtx"]
            adv = {}
            for cp, name in cmap.items():
                try:
                    adv[chr(cp)] = hmtx[name][0] / upem * spec[1]
                except Exception:
                    pass
            tt.close()
            if adv:
                entry = adv
        except Exception:
            entry = None
    _metrics_cache[key] = entry
    return entry


def str_w(s, size, font=FONT, bold=False, stress=1.0):
    """Width of a string in points."""
    adv = _metrics(font, bold)
    if adv is None:
        f = FALLBACK_CW.get(font, 0.50) * (1.05 if bold else 1.0)
        return len(s) * f * size * stress
    dflt = adv.get("n", 0.5)
    return sum(adv.get(c, dflt) for c in s) * size * stress


def wrap_lines(runs, size, w_in, stress=1.0):
    """Greedy word wrap over a list of runs; returns the line count."""
    limit = w_in * 72.0
    tokens = []          # (text, bold, italic, font)
    for r in runs:
        parts = r["text"].split(" ")
        for i, part in enumerate(parts):
            if i:
                tokens.append((" ", r))
            if part:
                tokens.append((part, r))
    lines, cur = 1, 0.0
    for tok, r in tokens:
        tw = str_w(tok, size, r.get("font", FONT), r.get("bold", False), stress)
        if tok == " ":
            cur += tw
            continue
        if cur + tw > limit + 0.01 and cur > 0:
            lines += 1
            cur = tw
        else:
            cur += tw
    return lines


LAYOUT_STRESS = 1.04   # size every box as if the real face were 4% wider


def est_lines(paras, w, stress=LAYOUT_STRESS):
    """Height in inches needed by a list of paragraphs in a box `w` wide."""
    total = 0.0
    for p in paras:
        n = wrap_lines(p["runs"], p["size"], w, stress)
        total += n * p["size"] * 1.21 * p.get("line", LINE)
        total += p.get("space_after", 0.0)
    return total / 72.0


shapes_log = []   # (slide_idx, kind, x, y, w, h, label)
text_log = []     # (slide_idx, label, w, h, paras)


def line_h(size, line=LINE):
    return size * 1.21 * line / 72.0


# ------------------------------------------------------------------- shapes --
def rect(slide, si, x, y, w, h, fill=None, line=None, line_w=1.0,
         shape=MSO_SHAPE.RECTANGLE, dash=False, adj=None, kind="rect"):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if adj is not None and sh.adjustments:
        sh.adjustments[0] = adj
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
        if dash:
            from lxml import etree
            ln = sh.line._get_or_add_ln()
            d = etree.SubElement(ln, qn("a:prstDash"))
            d.set("val", "dash")
    sh.text_frame.word_wrap = True
    shapes_log.append((si, kind, x, y, w, h, kind))
    return sh


def text(slide, si, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP,
         align=PP_ALIGN.LEFT, label="", kind="text", check=True):
    """paras: [{"size":.., "line":.., "space_after":.., "align":..,
                "runs":[{"text":..,"bold":..,"italic":..,"color":..,"font":..}]}]"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", align)
        para.line_spacing = p.get("line", LINE)
        if p.get("space_after"):
            para.space_after = Pt(p["space_after"])
        if p.get("hang"):
            pPr = para._p.get_or_add_pPr()
            pPr.set("marL", str(int(p["hang"] * 914400)))
            pPr.set("indent", str(-int(p["hang"] * 914400)))
        for r in p["runs"]:
            run = para.add_run()
            run.text = r["text"]
            f = run.font
            f.name = r.get("font", FONT)
            f.size = Pt(p["size"])
            f.bold = r.get("bold", False)
            f.italic = r.get("italic", False)
            f.color.rgb = r.get("color", INK)
            if r.get("spc"):
                run.font._rPr.set("spc", str(int(r["spc"] * 100)))
    shapes_log.append((si, kind, x, y, w, h, label or kind))
    if check:
        eff_w = w - max((q.get("hang", 0.0) for q in paras), default=0.0)
        text_log.append((si, label or p["runs"][0]["text"][:30], eff_w, h, paras))
    return tb


def para(txt, size=BODY_SZ, bold=False, italic=False, color=INK, font=FONT,
         line=LINE, space_after=0.0, align=None, label=None, spc=None,
         hang=0.0):
    """One paragraph, optionally with a bold lead-in label."""
    runs = []
    if label:
        runs.append({"text": label, "bold": True, "color": color, "font": font})
    if txt:
        runs.append({"text": txt, "bold": bold, "italic": italic,
                     "color": color, "font": font, "spc": spc})
    p = {"size": size, "line": line, "space_after": space_after, "runs": runs}
    if hang:
        p["hang"] = hang
    if align is not None:
        p["align"] = align
    return p


# ------------------------------------------------------------- page furniture --
def slide_title(slide, si, sentence, page, kicker=None, y=0.44,
                size=22.0, w=None):
    """Left-aligned descriptive title with a short thin rule to its left."""
    ty = y
    if kicker:
        text(slide, si, TITLE_X, ty, 11.0, 0.26,
             [para(kicker, size=FURNITURE_SZ, bold=True, color=ACCENT, spc=1.4)],
             label="kicker")
        ty += 0.32
    tw = w or (SW - TITLE_X - 1.30)
    paras = [para(sentence, size=size, bold=True, color=NAVY, line=0.98)]
    h = est_lines(paras, tw)
    text(slide, si, TITLE_X, ty, tw, h + 0.02, paras, label="title")
    # rule sits at the vertical centre of the title's first line
    rect(slide, si, RULE_X, ty + line_h(size, 0.98) / 2 - RULE_H / 2,
         RULE_W, RULE_H, fill=ACCENT, kind="rule")
    if page:
        text(slide, si, SW - 1.30, SH - 0.50, 0.68, 0.26,
             [para(page, size=FURNITURE_SZ, color=SLATE, align=PP_ALIGN.RIGHT)],
             align=PP_ALIGN.RIGHT, label="page", check=False)
    return ty + h + 0.02


# ---------------------------------------------------------------- the block --
def block(slide, si, x, y, w, label, items, banner_fill=NAVY):
    """Banner + dot-rail bullet block. Returns the y cursor after the block.

    items: list of dicts
        {"lvl":0, "label":"Bold lead-in:", "text":" rest of the sentence."}
        {"lvl":1, "text":"sub-bullet"}
        {"lvl":2, "text":"third level"}          -> rendered with an en dash
        optional: "italic", "font", "size"
    """
    rect(slide, si, x, y, w, BANNER_H, fill=banner_fill,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.20, kind="banner")
    text(slide, si, x + LABEL_DX, y + (BANNER_H - line_h(LABEL_SZ, 1.0)) / 2,
         w - LABEL_DX - 0.20, line_h(LABEL_SZ, 1.0) + 0.02,
         [para(label, size=LABEL_SZ, bold=True, color=WHITE, line=1.0,
               spc=0.6)],
         label=f"banner:{label}", kind="banner_label")

    cy = y + BANNER_H + GAP_BANNER
    rail_top = cy - 0.05
    donut_ys = []

    for it in items:
        lvl = it.get("lvl", 0)
        if lvl == 0:
            size = it.get("size", BODY_SZ)
            tx, tw = x + TEXT_DX, w - TEXT_DX - 0.08
            color = INK
        elif lvl == 1:
            size = it.get("size", SUB_SZ)
            tx, tw = x + SUB_TEXT_DX, w - SUB_TEXT_DX - 0.08
            color = INK_SOFT
        else:
            size = it.get("size", L3_SZ)
            tx, tw = x + L3_TEXT_DX, w - L3_TEXT_DX - 0.08
            color = INK_SOFT

        body = it.get("text", "")
        if lvl == 2:
            body = "– " + body.lstrip()
        hang = it.get("hang", 0.0)
        paras = [para(body, size=size, italic=it.get("italic", False),
                      color=color, font=it.get("font", FONT),
                      label=it.get("label"), hang=hang)]
        h = est_lines(paras, tw - hang)
        text(slide, si, tx, cy, tw, h + 0.02, paras,
             label=f"{label}/L{lvl}/{(it.get('label') or body)[:26]}")

        first = line_h(size)
        if lvl == 0:
            ccy = cy + first / 2
            donut_ys.append(ccy)
        elif lvl == 1:
            rect(slide, si, x + SUB_MARK_DX, cy + first / 2 - SQ / 2, SQ, SQ,
                 fill=ACCENT, kind="square")
        cy += h + 0.02 + (GAP_ITEM if lvl == 0 else GAP_SUB)

    # rail behind the donuts: drawn now, then donuts on top of it
    if donut_ys:
        rail_bot = donut_ys[-1]
        rect(slide, si, x + RAIL_DX - RAIL_W / 2, rail_top, RAIL_W,
             max(0.02, rail_bot - rail_top), fill=RAIL, kind="rail")
        for ccy in donut_ys:
            rect(slide, si, x + RAIL_DX - DONUT_D / 2, ccy - DONUT_D / 2,
                 DONUT_D, DONUT_D, fill=WHITE, line=ACCENT, line_w=1.5,
                 shape=MSO_SHAPE.OVAL, kind="donut")

    return cy - (GAP_ITEM if items and items[-1].get("lvl", 0) == 0
                 else GAP_SUB) + GAP_BLOCK


# =============================================================================
prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
blank = prs.slide_layouts[6]

L_X, L_W = M_L, 7.44                     # main column
R_X, R_W = 8.42, M_L + CONTENT_W - 8.42  # narrower second column (4.29)

# ------------------------------------------------- SLIDE 1 — title / problem --
s1 = prs.slides.add_slide(blank)
si = 1
text(s1, si, TITLE_X, 0.44, 6.0, 0.26,
     [para("ABERDEEN HACKATHON  ·  TEAM 05", size=FURNITURE_SZ, bold=True,
           color=ACCENT, spc=1.4)], label="s1 kicker")
_p = [para("RFP Pursuit Copilot", size=34.0, bold=True, color=NAVY, line=0.98)]
text(s1, si, TITLE_X, 0.74, 8.0, line_h(34.0, 0.98) + 0.02, _p, label="s1 name")
rect(s1, si, RULE_X, 0.74 + line_h(34.0, 0.98) / 2 - RULE_H / 2,
     RULE_W, RULE_H, fill=ACCENT, kind="rule")
_p = [para("Turns a 30–80 page RFP into win themes, a pursuit strategy and a "
           "first-draft proposal, for the Aberdeen pursuit team.",
           size=14.0, color=SLATE, line=1.10)]
text(s1, si, TITLE_X, 1.36, 10.6, est_lines(_p, 10.6) + 0.02, _p,
     label="s1 oneliner")
text(s1, si, SW - 1.30, SH - 0.50, 0.68, 0.26,
     [para("1 / 4", size=FURNITURE_SZ, color=SLATE, align=PP_ALIGN.RIGHT)],
     align=PP_ALIGN.RIGHT, label="s1 page", check=False)

y = 2.62
y = block(s1, si, L_X, y, L_W, "The Problem", [
    {"lvl": 0, "size": 14.0, "label": "Seven days to respond.",
     "text": " The first 48 hours go to structure and boilerplate, not to the "
             "client."},
    {"lvl": 0, "size": 14.0, "label": "Every bidder sounds the same.",
     "text": " Rivals send similar credentials and the same generic approach."},
    {"lvl": 0, "size": 14.0, "label": "The scoring punishes it.",
     "text": " Clients grade demonstrated understanding above capability lists."},
])
y = block(s1, si, L_X, y, L_W, "Who It's For", [
    {"lvl": 0, "size": 14.0,
     "text": "The BD lead, partner and SME sharing hours 0–48 of a "
             "seven-day window."},
])
S1L = y

ry = block(s1, si, R_X, 2.62, R_W, "What The Client Wrote Down", [
    {"lvl": 0, "font": QUOTE_FONT, "italic": True, "size": 12.5,
     "text": "“Wayfarer Market Co. puts 30% of its score on Understanding "
             "of Brand & Crew Culture — and explicitly rejects the "
             "cut-costs, loyalty-program, more-ads playbook.”"},
    {"lvl": 1, "label": "Read:",
     "text": " the generic proposal does not just underwhelm. It loses on the "
             "criteria the client wrote down."},
])
S1R = ry

text(s1, si, M_L, SH - 0.50, 6.4, 0.26,
     [para("Carrie Stout  ·  CJ Johnson  ·  Jordan Cook  ·  "
           "Preetish Rath", size=FURNITURE_SZ, color=SLATE)], label="s1 team")

# ------------------------------------------------------- SLIDE 2 — solution --
s2 = prs.slides.add_slide(blank)
si = 2
y = slide_title(s2, si, "Upload the RFP; the copilot returns a point of view, "
                        "not a template.", "2 / 4",
                kicker="THE SOLUTION", w=11.0) + 0.26
TOP2 = y

y = block(s2, si, L_X, y, L_W, "How It Works", [
    {"lvl": 0, "label": "Ingest.",
     "text": " RFP, Aberdeen credentials, prior proposals, case studies."},
    {"lvl": 0, "label": "Analyze.",
     "text": " Objectives, requirements, constraints, evaluation criteria."},
    {"lvl": 0, "label": "Draft.",
     "text": " Win themes first — then approach, staffing, timeline, copy."},
    {"lvl": 0, "label": "Score.",
     "text": " Predicted rubric score and completeness check before send."},
])
y = block(s2, si, L_X, y, L_W, "What It Hands You", [
    {"lvl": 0, "label": "Three win themes,",
     "text": " tied to the client's stated priorities."},
    {"lvl": 0, "label": "Solution approach,",
     "text": " staffing model and timeline."},
    {"lvl": 0, "label": "Aberdeen experience",
     "text": " ranked by relevance — and why it is relevant."},
    {"lvl": 0, "label": "First-draft copy",
     "text": " for the opening sections, structured to the client's own required "
             "response items."},
    {"lvl": 0, "label": "A Response Action per requirement:",
     "text": " what it asks of us."},
    {"lvl": 1, "text": "Controlled vocabulary, five values:"},
    {"lvl": 2, "text": "Address · Provide Information · Provide "
                       "Attachment · Acknowledge / Confirm · "
                       "Deliverable if Awarded"},
])
S2L = y

ry = block(s2, si, R_X, TOP2, R_W, "The Part Nobody Else Built", [
    {"lvl": 0, "label": "It grades itself",
     "text": " against the client's own weighted rubric before you send it."},
    {"lvl": 0, "text": "We have been the buy-side judge for years, so we already "
                       "hold the scoring sheets."},
    {"lvl": 0, "label": "Completeness check:",
     "text": " an unanswered requirement is a non-responsive bid."},
    {"lvl": 0, "label": "“Why Aberdeen”",
     "text": " comes from our Culture Charter — low-ego, high-ownership "
             "partnership — with every claim citing its source."},
])
S2R = ry

_p = [para("Drafts are written to sound like people wrote them: warm, specific "
           "and about the client. Review before send, always.",
           size=12.0, italic=True, color=SLATE, line=1.08)]
text(s2, si, M_L, 6.94, 11.0, est_lines(_p, 11.0) + 0.02, _p,
     label="s2 footer")

# ----------------------------------------------------------- SLIDE 3 — demo --
s3 = prs.slides.add_slide(blank)
si = 3
y = slide_title(s3, si, "A Wayfarer Market Co. walkthrough, from upload to "
                        "predicted score.", "3 / 4", kicker="DEMO",
                w=11.0) + 0.26
top3 = y

y = block(s3, si, L_X, y, L_W, "The Input", [
    {"lvl": 0, "label": "Specialty grocer,",
     "text": " ~500 stores, ~50,000 crew. Grow revenue, keep crew culture "
             "intact."},
])
y = block(s3, si, L_X, y, L_W, "Walkthrough", [
    {"lvl": 0, "hang": 0.22, "label": "1  ",
     "text": "Copilot parses the RFP: scope, deliverables D1–D5, "
             "requirements 5.1–5.7, weighted criteria, page limit."},
    {"lvl": 0, "hang": 0.22, "label": "2  ",
     "text": "It pins the non-negotiable up front: nothing that erodes crew "
             "culture."},
    {"lvl": 0, "hang": 0.22, "label": "3  ",
     "text": "Win themes come back in Wayfarer's language — crew first, "
             "growth second."},
    {"lvl": 0, "hang": 0.22, "label": "4  ",
     "text": "The draft assembles against Exhibit B, the client's ten "
             "required response items."},
    {"lvl": 0, "hang": 0.22, "label": "5  ",
     "text": "Predicted score names the gap: the playbook fails the 30% "
             "culture criterion."},
])
S3L = y
_p = [para("What the judges see: a proposal that argues the client's own "
           "priorities back to them — and a number saying how it would be "
           "scored.", size=12.0, italic=True, color=SLATE, line=1.08)]
text(s3, si, L_X + TEXT_DX, y - GAP_BLOCK + 0.16, L_W - TEXT_DX - 0.08,
     est_lines(_p, L_W - TEXT_DX - 0.08) + 0.02, _p, label="s3 close")

# right column: banner + labelled screenshot placeholder
rect(s3, si, R_X, top3, R_W, BANNER_H, fill=NAVY,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.20, kind="banner")
text(s3, si, R_X + LABEL_DX,
     top3 + (BANNER_H - line_h(LABEL_SZ, 1.0)) / 2, R_W - LABEL_DX - 0.20,
     line_h(LABEL_SZ, 1.0) + 0.02,
     [para("The Screen", size=LABEL_SZ, bold=True, color=WHITE, line=1.0,
           spc=0.6)], label="banner:The Screen", kind="banner_label")
phy = top3 + BANNER_H + GAP_BANNER
phh = 3.90
rect(s3, si, R_X, phy, R_W, phh, fill=LIGHT, line=SLATE, line_w=1.25,
     dash=True, kind="placeholder")
_p = [para("[ Screenshot: paste UI capture here ]", size=13.0, bold=True,
           color=SLATE, align=PP_ALIGN.CENTER, space_after=7.0),
      para("Front end built; backend in progress at time of writing.",
           size=12.0, italic=True, color=SLATE, align=PP_ALIGN.CENTER)]
text(s3, si, R_X + 0.30, phy + phh / 2 - 0.42, R_W - 0.60, 0.84, _p,
     align=PP_ALIGN.CENTER, label="s3 placeholder", kind="ph_text")
text(s3, si, R_X, phy + phh + 0.14, R_W, 0.24,
     [para("Live walkthrough delivered in the demo slot.", size=FURNITURE_SZ,
           color=SLATE, align=PP_ALIGN.CENTER)], align=PP_ALIGN.CENTER,
     label="s3 ph caption")

# --------------------------------------- SLIDE 4 — impact / path to market ---
s4 = prs.slides.add_slide(blank)
si = 4
y = slide_title(s4, si, "Where the value sits, how we would pilot it, what could "
                        "go wrong.", "4 / 4",
                kicker="IMPACT  ·  PATH TO MARKET", w=11.0) + 0.26
top4 = y

y = block(s4, si, L_X, y, L_W, "Business Value", [
    {"lvl": 0, "text": "More pursuits answered, at steadier quality."},
    {"lvl": 0, "label": "Effort moves to the heaviest criteria:",
     "text": " demonstrated understanding is roughly 45% of a rubric, against "
             "~15% for functionality."},
])
y = block(s4, si, L_X, y, L_W, "Pilot Plan", [
    {"lvl": 0, "hang": 0.22, "label": "1  ",
     "text": "Run in parallel on the next live pursuit, beside the human draft."},
    {"lvl": 0, "hang": 0.22, "label": "2  ",
     "text": "Stand up a curated library of credentials, proposals and case "
             "studies."},
    {"lvl": 1, "text": "The main dependency — and empty today."},
    {"lvl": 0, "hang": 0.22, "label": "3  ",
     "text": "Partner review before anything ships. A rule, not a setting."},
])
y = block(s4, si, L_X, y, L_W, "Risks, Honestly", [
    {"lvl": 0, "text": "Invented credentials or client facts."},
    {"lvl": 0, "text": "One uniform voice across every pursuit."},
    {"lvl": 0, "text": "Confidentiality of client RFP material."},
    {"lvl": 0, "label": "Mitigation:",
     "text": " partner sign-off on anything sent, and a cited source behind "
             "every claim."},
])
S4L = y

ry = block(s4, si, R_X, top4, R_W, "What We'll Measure", [
    {"lvl": 0, "text": "Draft-to-first-review time."},
    {"lvl": 0, "text": "Rubric completeness rate."},
    {"lvl": 0, "text": "Predicted versus actual evaluator score."},
    {"lvl": 0, "text": "Win rate where it was used."},
    {"lvl": 1, "text": "Baselines set in the pilot. Nothing claimed yet."},
])
ry = block(s4, si, R_X, ry, R_W, "Aberdeen Labs Fit", [
    {"lvl": 0, "text": "Ingest-and-structure, score-against-a-rubric and "
                       "write-from-our-real-values lift straight into pitch "
                       "decks, QBRs and any pursuit artifact."},
    {"lvl": 0, "text": "The buy-side rubric library behind the scoring is an "
                       "Aberdeen asset no competitor can copy."},
])
S4R = ry

print(f"\ncolumn cursors  s1 L={S1L:.2f} R={S1R:.2f} | s2 L={S2L:.2f} R={S2R:.2f}"
      f" | s3 L={S3L:.2f} | s4 L={S4L:.2f} R={S4R:.2f}   (safe bottom 6.85)")

# ================================================================== save =====
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "hack-team-05-pursuit-copilot-deck.pptx")
prs.save(OUT)
print("saved", OUT)


# ============================================================ verification ===
def verify(path):
    p = Presentation(path)
    bad = 0
    print(f"\nslides: {len(p.slides._sldIdLst)}  "
          f"canvas: {p.slide_width / 914400:.3f} x {p.slide_height / 914400:.3f} in")
    if len(p.slides._sldIdLst) != 4:
        print("  !! slide count is not 4")
        bad += 1

    # Nominal must fit exactly; STRESS re-measures every string 4% wider to
    # cover the gap between our donor-font metrics and real Calibri/Georgia.
    STRESS = 1.08
    print("\n--- text fit: layout / +8% stress / box (inches) ---")
    tight = []
    for si, label, w, h, paras in text_log:
        est = est_lines(paras, w)
        stressed = est_lines(paras, w, STRESS)
        flag = ""
        if est > h + 0.004:
            flag = "  <== OVERFLOW"
            bad += 1
        elif stressed > h + 0.004:
            flag = "  <- tight at +8%"
            tight.append((si, label))
        print(f"s{si} {est:5.2f} /{stressed:5.2f} /{h:5.2f}  "
              f"{label[:48]}{flag}")
    print(f"  tight only under the advisory +8% stress: {len(tight)}")
    for si, label in tight:
        print(f"    s{si} {label[:60]}")

    print("\n--- geometry: off-slide ---")
    off = 0
    for si, kind, x, y, w, h, _lbl in shapes_log:
        if x < -0.001 or y < -0.001 or x + w > SW + 0.001 or y + h > SH + 0.001:
            print(f"  !! s{si} {kind} at {x:.2f},{y:.2f} {w:.2f}x{h:.2f}")
            off += 1
    bad += off
    print("  none" if not off else "")

    print("\n--- vertical extent per slide (content bottom, limit 7.50) ---")
    for si in sorted({r[0] for r in shapes_log}):
        bot = max(r[3] + r[5] for r in shapes_log if r[0] == si)
        print(f"  s{si}: {bot:.2f}")

    print("\n--- geometry: unexpected overlaps ---")
    nest_ok = {"banner_label", "rail", "donut", "square", "ph_text", "rule"}
    by_slide = {}
    for rec in shapes_log:
        by_slide.setdefault(rec[0], []).append(rec)
    ov = 0
    for si, recs in by_slide.items():
        for i in range(len(recs)):
            for j in range(i + 1, len(recs)):
                _, k1, x1, y1, w1, h1, l1 = recs[i]
                _, k2, x2, y2, w2, h2, l2 = recs[j]
                if k1 in nest_ok or k2 in nest_ok:
                    continue
                dx = min(x1 + w1, x2 + w2) - max(x1, x2)
                dy = min(y1 + h1, y2 + h2) - max(y1, y2)
                if dx > 0.01 and dy > 0.01:
                    print(f"  !! s{si} [{l1[:34]}] x [{l2[:34]}]  "
                          f"overlap {dx:.2f}x{dy:.2f} in")
                    ov += 1
    bad += ov
    if not ov:
        print("  none")

    print("\n--- fonts / sizes / fills in the saved file ---")
    fonts, sizes, fills = set(), set(), set()
    for s in p.slides:
        for sh in s.shapes:
            if sh.shape_type is not None and sh.has_text_frame is False:
                pass
            try:
                if sh.fill.type is not None and sh.fill.type == 1:
                    fills.add(str(sh.fill.fore_color.rgb))
            except Exception:
                pass
            if sh.has_text_frame:
                for pa in sh.text_frame.paragraphs:
                    for r in pa.runs:
                        if r.font.name:
                            fonts.add(r.font.name)
                        if r.font.size:
                            sizes.add(round(r.font.size.pt, 1))
    print("  fonts:", sorted(fonts))
    print("  sizes:", sorted(sizes))
    print("  solid fills:", sorted(fills))
    small = [s for s in sizes if s < 12.0]
    if small:
        print("  !! sizes below 12pt:", small)
        bad += len(small)

    print(f"\nfile size: {os.path.getsize(path) / 1024:.1f} KB")
    print("PROBLEMS:", bad)
    return bad


verify(OUT)
